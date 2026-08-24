import sys
try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i + 1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
        print("\n")

extract_text_from_pptx("/Users/natthawutjantakul/intelligist_dataX/DEX_TOR_Gap_Recheck_20260624.pptx")
